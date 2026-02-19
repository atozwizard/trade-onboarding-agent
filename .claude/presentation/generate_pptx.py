"""
Trade Onboarding Agent — PPTX 자동 생성 스크립트
폰트: Noto Sans KR (시스템에 설치되어 있어야 함)
실행: uv run python .claude/presentation/generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─────────────────────────────────────────
# 색상 팔레트 (Upstage/SeSAC 스타일)
# ─────────────────────────────────────────
C_NAVY    = RGBColor(0x1A, 0x23, 0x3A)   # 배경/강조 (진남색)
C_BLUE    = RGBColor(0x2D, 0x6A, 0xE0)   # 메인 포인트 컬러
C_LIGHT   = RGBColor(0xF4, 0xF6, 0xFB)   # 슬라이드 배경 (연회색)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # 본문 텍스트
C_GRAY    = RGBColor(0x6B, 0x7B, 0x8D)   # 보조 텍스트
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)   # QuizAgent
C_ORANGE  = RGBColor(0xE6, 0x7E, 0x22)   # EmailAgent
C_RED     = RGBColor(0xC0, 0x39, 0x2B)   # RiskAgent

FONT = "Noto Sans KR"
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "TradeOnboardingAgent_발표.pptx")

# ─────────────────────────────────────────
# 헬퍼 함수
# ─────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs, bg_color=None):
    blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃
    slide = prs.slides.add_slide(blank_layout)
    if bg_color:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return slide


def add_rect(slide, x, y, w, h, color):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=16, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, x, y, w, h,
                  font_size=16, bold=False, color=C_DARK,
                  align=PP_ALIGN.LEFT, line_spacing=None):
    """줄바꿈 포함 멀티라인 텍스트박스"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    if line_spacing:
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line == "":
            p.add_run().text = ""
            continue
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


def section_badge(slide, number, label, bg=C_NAVY):
    """좌측 섹션 뱃지 (01. 팀 소개 등)"""
    add_rect(slide, 0, 0, 0.8, 7.5, bg)
    add_textbox(slide, f"0{number}", 0.05, 0.2, 0.7, 0.5,
                font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # 세로 텍스트는 pptx에서 직접 지원이 어려우므로 레이블은 생략하거나 회전 필요
    return


def slide_header(slide, title, subtitle=None):
    """슬라이드 상단 제목 영역"""
    add_rect(slide, 0.8, 0, 12.53, 1.0, C_NAVY)
    add_textbox(slide, title, 1.0, 0.2, 10.0, 0.7,
                font_size=20, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 1.0, 0.65, 10.0, 0.4,
                    font_size=12, color=RGBColor(0xAA, 0xBB, 0xCC))


# ─────────────────────────────────────────
# 슬라이드 01: 표지
# ─────────────────────────────────────────

def slide_cover(prs):
    slide = blank_slide(prs, C_NAVY)

    # 상단 포인트 바
    add_rect(slide, 0, 0, 13.33, 0.12, C_BLUE)

    # 메인 타이틀
    add_textbox(slide, "Trade Onboarding Agent",
                1.5, 1.8, 10.0, 1.2,
                font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 서브타이틀
    add_textbox(slide, "채팅으로 시작하는 무역 실무 온보딩",
                1.5, 3.0, 10.0, 0.7,
                font_size=22, bold=False, color=RGBColor(0xAA, 0xCC, 0xFF),
                align=PP_ALIGN.CENTER)

    # 팀 정보
    add_textbox(slide, "이성준  ·  차세종  ·  황지은  ·  이영기",
                1.5, 4.0, 10.0, 0.5,
                font_size=16, color=RGBColor(0xCC, 0xDD, 0xEE),
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "발표자: 이성준",
                1.5, 4.5, 10.0, 0.4,
                font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                align=PP_ALIGN.CENTER)

    # 하단
    add_rect(slide, 0, 7.1, 13.33, 0.4, C_BLUE)
    add_textbox(slide, "© 2026 Upstage Co., Ltd.   |   SeSAC AIPE Project",
                0, 7.15, 13.33, 0.3,
                font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 📌 이미지 안내 (실제 삽입 시 대체)
    add_textbox(slide, "[ 📌 권장: 항구·컨테이너 배경 이미지 삽입 ]",
                1.5, 5.2, 10.0, 0.4,
                font_size=11, color=RGBColor(0x66, 0x77, 0x88),
                align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────
# 슬라이드 02: 팀 소개
# ─────────────────────────────────────────

def slide_team(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 1, "팀소개")
    slide_header(slide, "팀 소개 및 역할", "01. 팀 소개 및 주제")

    members = [
        ("이성준", "QuizAgent · EvalTool 개발", C_GREEN),
        ("차세종", "RiskManagingAgent 개발", C_RED),
        ("황지은", "EmailAgent 개발", C_ORANGE),
        ("이영기", "프론트엔드 · RAG 파이프라인", C_BLUE),
    ]

    positions = [(1.0, 1.3), (4.0, 1.3), (7.0, 1.3), (10.0, 1.3)]

    for (name, role, color), (px, py) in zip(members, positions):
        add_rect(slide, px, py, 2.8, 3.5, color)
        add_textbox(slide, name, px, py + 0.3, 2.8, 0.7,
                    font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, role, px + 0.1, py + 1.1, 2.6, 1.5,
                    font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "* 역할 분담은 최종 확정 후 업데이트 예정",
                1.0, 5.1, 11.0, 0.4,
                font_size=10, color=C_GRAY)


# ─────────────────────────────────────────
# 슬라이드 03: 문제 정의
# ─────────────────────────────────────────

def slide_problem(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 1, "주제")
    slide_header(slide, "왜 무역 신입사원 온보딩이 어려운가?", "01. 팀 소개 및 주제")

    bullets = [
        "● 무역 실무 용어·프로세스는 암기보다 반복 경험이 필요",
        "● 기존 온보딩: 문서 중심 → 단방향 학습 → 낮은 몰입도",
        "● 실수 한 번이 페널티·클레임·계약 분쟁으로 직결되는 고위험 도메인",
        "● AI 코치가 대화하며 실시간 피드백을 주는 시스템 부재",
    ]
    add_multiline(slide, bullets, 1.0, 1.3, 11.0, 3.5, font_size=16, color=C_DARK)

    # 핵심 메시지 강조 박스
    add_rect(slide, 1.0, 5.0, 11.3, 1.0, C_NAVY)
    add_textbox(slide, '"실무는 교과서가 아니라 대화로 배운다"',
                1.2, 5.15, 11.0, 0.7,
                font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "📌 권장: 문서 더미 vs 채팅 버블 Before/After 이미지 삽입",
                1.0, 1.2, 11.0, 0.3, font_size=10, color=C_GRAY)


# ─────────────────────────────────────────
# 슬라이드 04: 서비스 기획
# ─────────────────────────────────────────

def slide_planning(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 2, "기획")
    slide_header(slide, "도메인 선정 및 시장 분석", "02. 서비스 기획 및 서비스 디자인")

    lines = [
        "● 타겟: 무역회사 신입사원 (1~2년차)",
        "",
        "● 핵심 학습 영역",
        "    ○ Incoterms / 결제 조건 (L/C, T/T)",
        "    ○ 무역 이메일 작성 및 검토",
        "    ○ 공급망 리스크 대응",
        "",
        "● 데이터 규모",
        "    ○ 17개 JSON 데이터셋  /  813 raw records  /  782 ingestable",
        "    ○ ICC 무역용어집 284개  +  화성상공회의소 용어 169개",
        "    ○ 실무 이메일·실수 사례·클레임 사례 포함",
    ]
    add_multiline(slide, lines, 1.0, 1.2, 11.5, 5.5, font_size=15, color=C_DARK)


# ─────────────────────────────────────────
# 슬라이드 05: 서비스 디자인 (시스템 흐름)
# ─────────────────────────────────────────

def slide_design(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 2, "디자인")
    slide_header(slide, "서비스 전체 흐름", "02. 서비스 기획 및 서비스 디자인")

    # 다이어그램 안내 박스
    add_rect(slide, 1.0, 1.2, 11.3, 4.5, RGBColor(0xE8, 0xEF, 0xFD))
    add_textbox(slide,
        "📌 기획서 시스템 개요 다이어그램 (graph LR) 캡처 이미지 삽입\n\n"
        "👤 신입사원 → 💬 React UI → 🎯 Orchestrator\n"
        "    → 📝 QuizAgent  (퀴즈 학습)\n"
        "    → ✉️  EmailAgent  (이메일 코칭)\n"
        "    → 🚨 RiskManagingAgent  (리스크 분석)",
        1.2, 1.4, 11.0, 4.0,
        font_size=15, color=C_DARK)

    add_multiline(slide, [
        "* Orchestrator가 사용자 의도를 자동 분류 → 전문 에이전트 라우팅",
        "* 단일 채팅창에서 모든 기능 이용 가능",
    ], 1.0, 5.9, 11.0, 0.8, font_size=12, color=C_GRAY)


# ─────────────────────────────────────────
# 슬라이드 06: Workflow 기획
# ─────────────────────────────────────────

def slide_workflow(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 3, "워크플로")
    slide_header(slide, "Agent Workflow 구성 — Orchestrator 중심 설계",
                 "03. Agent Workflow 기획 및 구성")

    lines_left = [
        "● LLM만으로 풀 수 없는 문제",
        "    ○ 무역 도메인 특화 지식 → RAG 필요",
        "    ○ 퀴즈 품질 보장 → EvalTool 검증 필요",
        "    ○ 리스크 분석 → 멀티턴 + 구조화 보고서 필요",
        "",
        "● Orchestrator 라우팅 우선순위",
        "    ① active_agent 유지 (멀티턴 진행 중)",
        "    ② context.mode 명시 (프론트 오버라이드)",
        "    ③ LLM 인텐트 분류 (solar-pro2)",
        "    ④ DefaultChatAgent (폴백)",
    ]
    add_multiline(slide, lines_left, 1.0, 1.2, 6.0, 5.5, font_size=14, color=C_DARK)

    # 오른쪽: 다이어그램 안내
    add_rect(slide, 7.2, 1.2, 5.9, 5.5, RGBColor(0xE8, 0xEF, 0xFD))
    add_textbox(slide,
        "📌 기획서 Orchestrator\n라우팅 flowchart TD\n캡처 이미지 삽입",
        7.4, 2.5, 5.5, 3.0,
        font_size=14, color=C_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────
# 슬라이드 07: RAG + Tool
# ─────────────────────────────────────────

def slide_rag(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 3, "RAG")
    slide_header(slide, "RAG 파이프라인 및 Tool 활용",
                 "03. Agent Workflow 기획 및 구성")

    lines = [
        "● RAG (Retrieval-Augmented Generation)",
        "    ○ 벡터 DB: ChromaDB (영구 저장, cosine 유사도)",
        "    ○ 임베딩: Upstage Solar Embedding (1024차원)",
        "    ○ 17개 JSON → 782 ingestable records 임베딩",
        "",
        "● 외부/내부 Tool",
        "    ○ EvalTool: 퀴즈 품질 검증 (is_valid + issues 반환)",
        "    ○ TradeTermValidator: 무역 용어 정확성 검증",
        "    ○ UnitValidator: 무게/부피/컨테이너 단위 일관성 검증",
        "    ○ Upstage Document Parse API: ICC PDF → JSON 변환 (OCR)",
    ]
    add_multiline(slide, lines, 1.0, 1.2, 7.5, 4.5, font_size=14, color=C_DARK)

    # 오른쪽: 파이프라인 다이어그램
    add_rect(slide, 8.7, 1.2, 4.4, 4.5, RGBColor(0xE8, 0xEF, 0xFD))
    add_textbox(slide,
        "📌 기획서 데이터 파이프라인\nflowchart LR 캡처 이미지 삽입\n\n"
        "dataset/*.json\n→ ingest.py\n→ ChromaDB\n→ retriever.py\n→ 각 에이전트",
        8.9, 1.5, 4.0, 4.0,
        font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────
# 슬라이드 08: 3개 에이전트 상세
# ─────────────────────────────────────────

def slide_agents(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 3, "에이전트")
    slide_header(slide, "3개 전문 에이전트",
                 "03. Agent Workflow 기획 및 구성")

    agents = [
        {
            "icon": "📝",
            "name": "QuizAgent",
            "color": C_GREEN,
            "lines": [
                "RAG 기반 퀴즈 5문제 생성",
                "EvalTool 5항목 자동 검증",
                "재시도 + 대체 생성 루프",
                "난이도: easy / medium / hard",
                "ICC 용어집 284개 활용",
            ]
        },
        {
            "icon": "✉️",
            "name": "EmailAgent",
            "color": C_ORANGE,
            "lines": [
                "Draft / Review 2가지 모드",
                "이메일 코칭 + 리스크 탐지",
                "무역 용어 오류 자동 검출",
                "단위 일관성 검증",
                "7개 서비스 병렬 검증",
            ]
        },
        {
            "icon": "🚨",
            "name": "RiskAgent",
            "color": C_RED,
            "lines": [
                "멀티턴 대화로 정보 수집",
                "RAG 유사 사례 검색",
                "5단계 리스크 평가",
                "JSON 구조화 보고서 생성",
                "예방 전략 단기/장기 제시",
            ]
        },
    ]

    col_x = [1.0, 4.9, 8.8]

    for agent, cx in zip(agents, col_x):
        # 헤더 박스
        add_rect(slide, cx, 1.2, 3.7, 0.7, agent["color"])
        add_textbox(slide, f"{agent['icon']}  {agent['name']}",
                    cx + 0.1, 1.3, 3.5, 0.5,
                    font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 본문 박스
        add_rect(slide, cx, 1.9, 3.7, 4.5, RGBColor(0xF8, 0xF9, 0xFF))
        add_multiline(slide,
                      ["● " + l for l in agent["lines"]],
                      cx + 0.15, 2.0, 3.4, 4.2,
                      font_size=13, color=C_DARK)


# ─────────────────────────────────────────
# 슬라이드 09: 기술 스택 + 역할 분담
# ─────────────────────────────────────────

def slide_stack(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 4, "개발")
    slide_header(slide, "기술 스택 및 역할 분담",
                 "04. Agent Workflow 개발/평가 및 시연")

    # 기술 스택 표 (좌)
    headers = ["레이어", "기술", "역할"]
    rows = [
        ["패키지 관리", "uv", "Python 패키지 매니저"],
        ["프론트엔드", "React 18 + Vite 5", "채팅 UI"],
        ["백엔드", "FastAPI + Python 3.11", "RESTful API 서버"],
        ["LLM", "Upstage Solar (solar-pro2)", "자연어 이해/생성"],
        ["임베딩", "Upstage Solar Embedding", "문서 벡터화 (1024차원)"],
        ["벡터 DB", "ChromaDB", "RAG 문서 검색"],
        ["트레이싱", "LangSmith", "에이전트 실행 디버깅"],
    ]

    y = 1.2
    add_rect(slide, 0.9, y, 8.2, 0.45, C_NAVY)
    for i, h in enumerate(headers):
        add_textbox(slide, h, 0.95 + i * 2.7, y + 0.05, 2.6, 0.35,
                    font_size=12, bold=True, color=C_WHITE)

    for r_idx, row in enumerate(rows):
        bg = RGBColor(0xEE, 0xF2, 0xFF) if r_idx % 2 == 0 else C_WHITE
        add_rect(slide, 0.9, y + 0.45 + r_idx * 0.42, 8.2, 0.42, bg)
        for c_idx, cell in enumerate(row):
            add_textbox(slide, cell,
                        0.95 + c_idx * 2.7, y + 0.48 + r_idx * 0.42,
                        2.6, 0.38, font_size=11, color=C_DARK)

    # 역할 분담 (우)
    add_rect(slide, 9.4, 1.2, 3.7, 0.45, C_NAVY)
    add_textbox(slide, "역할 분담", 9.5, 1.25, 3.5, 0.35,
                font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    roles = [
        ("이성준", "QuizAgent · EvalTool", C_GREEN),
        ("차세종", "RiskManagingAgent", C_RED),
        ("황지은", "EmailAgent", C_ORANGE),
        ("이영기", "프론트엔드 · RAG", C_BLUE),
    ]
    for i, (name, role, color) in enumerate(roles):
        add_rect(slide, 9.4, 1.65 + i * 1.1, 3.7, 1.0, color)
        add_textbox(slide, name, 9.5, 1.7 + i * 1.1, 3.5, 0.45,
                    font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, role, 9.5, 2.1 + i * 1.1, 3.5, 0.45,
                    font_size=12, color=C_WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────
# 슬라이드 10: 평가 방법 (EvalTool)
# ─────────────────────────────────────────

def slide_eval(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 4, "평가")
    slide_header(slide, "Agent Workflow 평가 — EvalTool 품질 검증",
                 "04. Agent Workflow 개발/평가 및 시연")

    lines = [
        "● QuizAgent EvalTool 검증 항목 (5개)",
        "    ○ 문제(question): RAG 원본 반영 여부",
        "    ○ 정답 선택지: 원본 데이터 일치 여부",
        "    ○ 오답 선택지: 실존 용어 기반 여부 (완전 허구 배제)",
        "    ○ 정답 인덱스: RAG 재검색으로 재확인",
        "    ○ 해설: 원본 내용 부합 여부",
        "",
        "● 재시도 및 대체 생성 루프 (MAX_RETRY = 2)",
        "    ① is_valid = false → issues 피드백 포함 재생성",
        "    ② 2회 소진 → 다른 용어로 대체 생성",
        "    ③ 합격 5문제 달성 시 즉시 종료",
        "",
        "● 예상 품질 지표",
        "    ○ 퀴즈 합격률 ~90%   |   무역 용어 검증 ~90%   |   단위 검증 ~95%",
    ]
    add_multiline(slide, lines, 1.0, 1.2, 11.5, 5.8, font_size=14, color=C_DARK)


# ─────────────────────────────────────────
# 슬라이드 11: 데모
# ─────────────────────────────────────────

def slide_demo(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 4, "데모")
    slide_header(slide, "서비스 데모",
                 "04. Agent Workflow 개발/평가 및 시연")

    demo_items = [
        ("📝 퀴즈 학습", "QuizAgent", C_GREEN,
         "퀴즈 5문제 생성 화면\n(스크린샷 삽입)"),
        ("✉️ 이메일 검토", "EmailAgent", C_ORANGE,
         "FOV→FOB 오류 검출 화면\n(스크린샷 삽입)"),
        ("🚨 리스크 분석", "RiskAgent", C_RED,
         "JSON 보고서 시각화\n(스크린샷 삽입)"),
    ]

    col_x = [1.0, 4.9, 8.8]

    for item, cx in zip(demo_items, col_x):
        icon_title, agent_name, color, placeholder = item
        add_rect(slide, cx, 1.2, 3.7, 0.6, color)
        add_textbox(slide, icon_title, cx + 0.1, 1.3, 3.5, 0.4,
                    font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, cx, 1.8, 3.7, 3.8, RGBColor(0xDD, 0xE8, 0xFF))
        add_textbox(slide, placeholder, cx + 0.2, 2.8, 3.3, 1.8,
                    font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 영상 링크
    add_rect(slide, 1.0, 5.8, 11.3, 0.7, C_NAVY)
    add_textbox(slide, "▶  전체 데모 영상: [ 영상 파일 삽입 또는 URL ]",
                1.2, 5.9, 11.0, 0.5,
                font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────
# 슬라이드 12: 마무리
# ─────────────────────────────────────────

def slide_conclusion(prs):
    slide = blank_slide(prs, C_LIGHT)
    section_badge(slide, 4, "마무리")
    slide_header(slide, "결론 및 향후 발전 방향",
                 "04. Agent Workflow 개발/평가 및 시연")

    lines_left = [
        "● 구현 성과",
        "    ○ RAG 기반 3개 전문 에이전트 완성",
        "    ○ EvalTool 자동 품질 검증 루프 구현",
        "    ○ 멀티턴 리스크 분석 + JSON 구조화 보고서",
        "    ○ 782개 무역 지식 ChromaDB 구축",
    ]
    lines_right = [
        "● 향후 발전 방향",
        "    ○ Redis 기반 세션 영속성 전환",
        "    ○ 사용자 학습 이력 추적 + 맞춤형 퀴즈",
        "    ○ 도메인 데이터 확충 (국가별 규정 등)",
        "    ○ 멀티모달: 계약서 PDF 업로드 → 분석",
    ]

    add_multiline(slide, lines_left, 1.0, 1.2, 5.8, 5.0, font_size=15, color=C_DARK)
    add_multiline(slide, lines_right, 7.0, 1.2, 5.8, 5.0, font_size=15, color=C_DARK)

    # 구분선
    add_rect(slide, 6.6, 1.2, 0.08, 5.0, C_BLUE)


# ─────────────────────────────────────────
# 슬라이드 13: Q&A
# ─────────────────────────────────────────

def slide_qa(prs):
    slide = blank_slide(prs, C_NAVY)
    add_rect(slide, 0, 0, 13.33, 0.12, C_BLUE)

    add_textbox(slide, "Q & A",
                0, 2.5, 13.33, 1.5,
                font_size=60, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "감사합니다.",
                0, 4.2, 13.33, 0.8,
                font_size=24, color=RGBColor(0xAA, 0xCC, 0xFF),
                align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 7.1, 13.33, 0.4, C_BLUE)
    add_textbox(slide, "© 2026 Upstage Co., Ltd.   |   SeSAC AIPE Project",
                0, 7.15, 13.33, 0.3,
                font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────

def main():
    prs = new_prs()

    print("Generating slides...")
    slide_cover(prs)       # 01 표지
    slide_team(prs)        # 02 팀 소개
    slide_problem(prs)     # 03 문제 정의
    slide_planning(prs)    # 04 서비스 기획
    slide_design(prs)      # 05 서비스 디자인
    slide_workflow(prs)    # 06 Workflow 기획
    slide_rag(prs)         # 07 RAG + Tool
    slide_agents(prs)      # 08 3개 에이전트
    slide_stack(prs)       # 09 기술 스택 + 역할 분담
    slide_eval(prs)        # 10 평가 방법
    slide_demo(prs)        # 11 데모
    slide_conclusion(prs)  # 12 마무리
    slide_qa(prs)          # 13 Q&A

    prs.save(OUTPUT_PATH)
    print("Done: " + OUTPUT_PATH)
    print("Total slides: " + str(len(prs.slides)))


if __name__ == "__main__":
    main()
